from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(16, 24, 39)

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # --- Slide 1: Problem & Opportunity ---
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1)
    add_title(slide1, "Problem & Opportunity")

    hex_w = Inches(1.5)
    hex_h = Inches(1.3)
    cx, cy = Inches(6.66), Inches(4)
    
    h_center = slide1.shapes.add_shape(MSO_SHAPE.HEXAGON, cx - hex_w/2, cy - hex_h/2, hex_w, hex_h)
    h_center.fill.solid()
    h_center.fill.fore_color.rgb = RGBColor(15, 23, 42)
    h_center.line.color.rgb = RGBColor(59, 130, 246)
    h_center.line.width = Pt(2)
    h_center.text_frame.text = "5W1H"
    h_center.text_frame.paragraphs[0].font.color.rgb = RGBColor(59, 130, 246)
    h_center.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    offsets = [(0, -1.2), (1.1, -0.6), (1.1, 0.6), (0, 1.2), (-1.1, 0.6), (-1.1, -0.6)]
    labels = ["What", "Why", "When", "How", "Where", "Who"]
    for i, (dx, dy) in enumerate(offsets):
        hx = slide1.shapes.add_shape(MSO_SHAPE.HEXAGON, cx + Inches(dx) - hex_w/2, cy + Inches(dy) - hex_h/2, hex_w, hex_h)
        hx.fill.solid()
        hx.fill.fore_color.rgb = RGBColor(59, 130, 246)
        hx.line.color.rgb = RGBColor(59, 130, 246)
        hx.text_frame.text = labels[i]
        hx.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        hx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_side_box(slide, left, top, width, height, text):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        box.line.color.rgb = RGBColor(59, 130, 246)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(203, 213, 225)
    
    add_side_box(slide1, 0.5, 1.8, 3.5, 1.4, "Manual exploratory analysis is slow and fragmented, causing delayed insights and limited ability to scale analytics work.")
    add_side_box(slide1, 0.5, 3.4, 3.5, 1.4, "Key users include business analysts, operations leads, and SME decision makers who need quick, reliable data-driven answers.")
    add_side_box(slide1, 0.5, 5.0, 3.5, 1.4, "These challenges appear across data-driven organizations facing complex, distributed datasets without massive data science teams.")
    
    add_side_box(slide1, 9.3, 1.8, 3.5, 1.4, "Because teams rely on scarce experts and manual coding, analysis is slow, error risk is higher, and decisions lack timely data support.")
    add_side_box(slide1, 9.3, 3.4, 3.5, 1.4, "Today, time-to-analysis stretches from hours to weeks; teams feel this delay during urgent incident investigations and reporting.")
    add_side_box(slide1, 9.3, 5.0, 3.5, 1.4, "An AI data analyst agent automates EDA, offers a self-service interface, and scales insight generation through natural language.")

    # --- Slide 2: System Capabilities ---
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2)
    add_title(slide2, "System Capabilities at a Glance")

    caps = [
        ("Core Data & Analysis\nFeatures", "Upload CSV/Excel files, run automated cleaning and exploratory analysis, generate charts, and transform messy raw data into structured insights."),
        ("Advanced Intelligence &\nInsights", "Perform anomaly detection, support natural language querying, generate concise insight reports, and guide users to defensible business decisions."),
        ("User Journey & Typical\nUse Cases", "Follow a clear path from raw files through exploration to decisions, supporting business analytics and reporting scenarios across all teams."),
        ("No-Code Experience &\nData Support", "Provide a fully no-code user experience with low-friction onboarding, accepting standard tabular formats for all complexity levels.")
    ]

    for i, (title, desc) in enumerate(caps):
        top = 1.8 + (i * 1.3)
        chev = slide2.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(0.5), Inches(top), Inches(3.5), Inches(1))
        chev.fill.solid()
        chev.fill.fore_color.rgb = RGBColor(37, 99, 235)
        chev.line.fill.background()
        p = chev.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        
        rect = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.8), Inches(top), Inches(9.0), Inches(1))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(30, 41, 59)
        rect.line.fill.background()
        p = rect.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(203, 213, 225)

    # --- Slide 3: End-to-End User Workflow ---
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3)
    add_title(slide3, "End-to-End User Workflow")

    div = slide3.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(6.5), Inches(3.5), Inches(0.33), Inches(0.5))
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(168, 85, 247)
    div.line.fill.background()

    l_line = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(0.05), Inches(3.5))
    l_line.fill.solid()
    l_line.fill.fore_color.rgb = RGBColor(59, 130, 246)
    l_line.line.fill.background()

    c1 = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(2.3), Inches(0.6), Inches(0.6))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(15, 23, 42)
    c1.line.color.rgb = RGBColor(59, 130, 246)
    c1.line.width = Pt(2)
    c1.text_frame.text = "01"

    t1 = slide3.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(4.5), Inches(2))
    t1.text_frame.word_wrap = True
    p = t1.text_frame.paragraphs[0]
    p.text = "Data Ingestion & Preparation"
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p = t1.text_frame.add_paragraph()
    p.text = "Users upload CSV or Excel files. The system automatically infers schema, profiles data quality, generates summary statistics, and embeds the metadata into a ChromaDB Vector Store for robust downstream analysis."
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(203, 213, 225)
    p.alignment = PP_ALIGN.CENTER

    r_line = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.3), Inches(2.5), Inches(0.05), Inches(3.5))
    r_line.fill.solid()
    r_line.fill.fore_color.rgb = RGBColor(59, 130, 246)
    r_line.line.fill.background()

    c2 = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(2.3), Inches(0.6), Inches(0.6))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(15, 23, 42)
    c2.line.color.rgb = RGBColor(168, 85, 247)
    c2.line.width = Pt(2)
    c2.text_frame.text = "02"

    t2 = slide3.shapes.add_textbox(Inches(7.3), Inches(3.2), Inches(4.5), Inches(2))
    t2.text_frame.word_wrap = True
    p = t2.text_frame.paragraphs[0]
    p.text = "Analysis, Exploration & Insight"
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p = t2.text_frame.add_paragraph()
    p.text = "The Smart Query Router triages questions. Simple queries run instantly on DuckDB. Complex queries trigger the LangChain agent for deep Python-based reasoning. Users receive fast, accurate answers and interactive Plotly charts."
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(203, 213, 225)
    p.alignment = PP_ALIGN.CENTER

    # --- Slide 4: Architecture Overview ---
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4)
    add_title(slide4, "System Architecture Overview")

    def add_arch_card(slide, left, top, num, title, desc, color):
        img_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.8), Inches(3))
        img_rect.fill.solid()
        img_rect.fill.fore_color.rgb = RGBColor(30, 41, 59)
        img_rect.line.fill.background()
        
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.3), Inches(top + 2.7), Inches(0.6), Inches(0.6))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        badge.text_frame.text = num

        tb = slide.shapes.add_textbox(Inches(left), Inches(top + 3.4), Inches(5.8), Inches(2))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p = tb.text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(203, 213, 225)
        p.space_before = Pt(5)

    add_arch_card(slide4, 0.5, 1.8, "01", "End-to-End Architecture & Data Flow", "Built on FastAPI and Next.js. Data flows from upload through automatic parsing, statistical profiling, and vector embedding into ChromaDB, separating data storage from the AI reasoning engine.", RGBColor(59, 130, 246))
    add_arch_card(slide4, 7.0, 1.8, "02", "Core Components & AI Routing", "Features a Hybrid Smart Query Router. Fast Track uses DuckDB for <3s latency. Complex Track uses a ReAct LangChain agent with Mistral fallbacks to ensure robust code execution.", RGBColor(168, 85, 247))

    # --- Slide 5: Data Ingestion Pipeline ---
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5)
    add_title(slide5, "Data Ingestion & Cleaning Pipeline")

    line = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(0.02), Inches(3.5))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(51, 65, 85)
    line.line.fill.background()

    def add_timeline_item(slide, top, title, desc):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), Inches(top), Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(59, 130, 246)
        dot.line.fill.background()
        idot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(top+0.1), Inches(0.1), Inches(0.1))
        idot.fill.solid()
        idot.fill.fore_color.rgb = RGBColor(255, 255, 255)
        idot.line.fill.background()
        
        tb1 = slide.shapes.add_textbox(Inches(1.5), Inches(top - 0.1), Inches(3.5), Inches(1))
        tb1.text_frame.word_wrap = True
        tb1.text_frame.paragraphs[0].text = title
        tb1.text_frame.paragraphs[0].font.size = Pt(18)
        tb1.text_frame.paragraphs[0].font.bold = True

        tb2 = slide.shapes.add_textbox(Inches(5.0), Inches(top - 0.1), Inches(7.5), Inches(1.5))
        tb2.text_frame.word_wrap = True
        tb2.text_frame.paragraphs[0].text = desc
        tb2.text_frame.paragraphs[0].font.size = Pt(14)
        tb2.text_frame.paragraphs[0].font.color.rgb = RGBColor(203, 213, 225)

    add_timeline_item(slide5, 2.5, "Automated parsing and preprocessing steps", "Parse CSV or Excel files, infer schemas and data types, analyze data bounds, and pre-compute minimums, maximums, and means for every column instantly upon upload.")
    add_timeline_item(slide5, 5.0, "Quality reporting and semantic memory", "Generate a statistical summary and inject it alongside table schemas directly into ChromaDB. This semantic memory layer allows the LLM to understand the dataset without re-reading the entire file.")

    # --- Slide 6: Anomaly Detection (Fishbone/Arrow) ---
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6)
    add_title(slide6, "Smart Anomaly Detection & Reasoning")

    arrow_line = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(9.5), Inches(0.15))
    arrow_line.fill.solid()
    arrow_line.fill.fore_color.rgb = RGBColor(59, 130, 246)
    arrow_line.line.fill.background()

    arrow_head = slide6.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(10.5), Inches(3.3), Inches(1.5), Inches(1.55))
    arrow_head.fill.solid()
    arrow_head.fill.fore_color.rgb = RGBColor(37, 99, 235)
    arrow_head.line.fill.background()
    arrow_head.rotation = 90

    def add_fishbone_section(slide, x, title_top, title, desc_top, desc, is_bottom=False):
        shape = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(x), Inches(2.5 if not is_bottom else 4.15), Inches(1), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
        shape.line.color.rgb = RGBColor(168, 85, 247) if is_bottom else RGBColor(59, 130, 246)
        
        tb = slide.shapes.add_textbox(Inches(x - 0.5), Inches(title_top), Inches(2.5), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(168, 85, 247) if is_bottom else RGBColor(59, 130, 246)
        p.alignment = PP_ALIGN.CENTER

        tb2 = slide.shapes.add_textbox(Inches(x - 0.5), Inches(desc_top), Inches(2.5), Inches(2))
        tb2.text_frame.word_wrap = True
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(203, 213, 225)

    add_fishbone_section(slide6, 2.5, 1.5, "Anomaly Types", 2.0, "Clarifies the difference between simple statistical boundaries (DuckDB IQR) and deep context-driven anomalies (LangChain).")
    add_fishbone_section(slide6, 5.5, 1.5, "Methods", 2.0, "Applies IQR, Z-scores, and iterative Python scripting dynamically based on the dataset schema.")
    add_fishbone_section(slide6, 8.5, 1.5, "Dimensions", 2.0, "Evaluates anomalies across categorical segments and time windows to highlight exactly where abnormal behavior occurs.")
    add_fishbone_section(slide6, 3.5, 6.0, "Time Series", 4.5, "Generates interactive Plotly charts dynamically to visually flag unusual patterns over time while highlighting trends.", is_bottom=True)
    add_fishbone_section(slide6, 6.5, 6.0, "Causality & Feedback", 4.5, "ReAct loop investigates causes by writing testing scripts. Triple-redundancy fallback guarantees successful execution.", is_bottom=True)

    # --- Slide 7: Natural Language Insights (Concentric Circles) ---
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7)
    add_title(slide7, "Natural Language Insights")

    cx, cy = 3.5, 4.0
    colors = [RGBColor(30, 58, 138), RGBColor(37, 99, 235), RGBColor(59, 130, 246)]
    radii = [2.5, 1.8, 1.1]
    labels = ["WHAT?", "HOW?", "WHY?"]

    for i in range(3):
        r = radii[i]
        circle = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(r*2), Inches(r*2))
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors[i]
        circle.line.fill.background()
        
        tb = slide7.shapes.add_textbox(Inches(cx - 0.5), Inches(cy - r + 0.2), Inches(1), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = labels[i]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    def add_nl_block(slide, top, title1, title2, desc):
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(top), Inches(1.5), Inches(1.2))
        sq.fill.solid()
        sq.fill.fore_color.rgb = RGBColor(37, 99, 235)
        sq.line.fill.background()
        sq.text_frame.word_wrap = True
        p = sq.text_frame.paragraphs[0]
        p.text = title1
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.0), Inches(top), Inches(4.5), Inches(1.2))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(30, 41, 59)
        rect.line.fill.background()
        rect.text_frame.word_wrap = True
        p2 = rect.text_frame.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(203, 213, 225)

    add_nl_block(slide7, 2.0, "Phenomenon\n/Outcome", "(Execution)", "Empowers non-technical users to converse with data, dynamically generating Pandas or DuckDB SQL code replacing manual effort.")
    add_nl_block(slide7, 3.5, "Method\n/Action", "(Method)", "Uses a Hybrid Smart Query Router. Simple questions run instantly via SQL. Complex questions trigger multi-step reasoning.")
    add_nl_block(slide7, 5.0, "Mission\n/Idea", "(Purpose)", "Aims to make analytics universally accessible yet fast and accurate, combining schema awareness with safe-answer execution.")

    # --- Slide 8: Technical Challenges (Venn Diagram) ---
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8)
    add_title(slide8, "Technical Challenges & Solutions")

    c1 = slide8.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(2.5), Inches(3.5), Inches(3.5))
    c1.fill.background()
    c1.line.color.rgb = RGBColor(59, 130, 246)
    c1.line.width = Pt(4)
    tb1 = slide8.shapes.add_textbox(Inches(4.2), Inches(4.0), Inches(1.5), Inches(0.5))
    tb1.text_frame.paragraphs[0].text = "Challenges"
    tb1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    tb1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    c2 = slide8.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.8), Inches(2.5), Inches(3.5), Inches(3.5))
    c2.fill.background()
    c2.line.color.rgb = RGBColor(168, 85, 247)
    c2.line.width = Pt(4)
    tb2 = slide8.shapes.add_textbox(Inches(6.8), Inches(4.0), Inches(1.5), Inches(0.5))
    tb2.text_frame.paragraphs[0].text = "Solutions"
    tb2.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    tb2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    arr = slide8.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.8), Inches(4.1), Inches(0.7), Inches(0.3))
    arr.fill.solid()
    arr.fill.fore_color.rgb = RGBColor(255, 255, 255)
    arr.line.fill.background()

    def add_venn_label(slide, left, top, num, text, right_align=False):
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(0.4), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(37, 99, 235)
        badge.line.fill.background()
        badge.text_frame.text = num
        
        tb = slide.shapes.add_textbox(Inches(left - 3.2 if right_align else left + 0.6), Inches(top - 0.2), Inches(3.0), Inches(1.5))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(203, 213, 225)
        p.alignment = PP_ALIGN.RIGHT if right_align else PP_ALIGN.LEFT

    add_venn_label(slide8, 3.2, 2.6, "1", "High latency: Complex LangChain reasoning loops took 15 to 30+ seconds.", right_align=True)
    add_venn_label(slide8, 2.6, 4.0, "2", "Parsing errors: LangChain crashed when outputting dual Action/Thoughts.", right_align=True)
    add_venn_label(slide8, 3.2, 5.4, "3", "Context limits: Large datasets caused token limits and hallucinations.", right_align=True)

    add_venn_label(slide8, 9.2, 2.6, "1", "Hybrid Architecture: DuckDB fast track reduced simple query latency to <3 seconds.", right_align=False)
    add_venn_label(slide8, 9.8, 4.0, "2", "Custom Extraction: Built regex-based fallback parsers to rescue final answers.", right_align=False)
    add_venn_label(slide8, 9.2, 5.4, "3", "ChromaDB Pipeline: Pre-computes metadata on upload for lean context retrieval.", right_align=False)

    prs.save('AI_Data_Analyst_Project_Presentation.pptx')

if __name__ == '__main__':
    create_presentation()
